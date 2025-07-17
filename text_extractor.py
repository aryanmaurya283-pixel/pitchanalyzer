import os
import tempfile
import docx2txt
import PyPDF2
import pptx

def extract_text_from_pdf(file):
    text = ""
    try:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + " "
    except Exception:
        return ""
    return text.strip()

def extract_text_from_docx(file):
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
            tmp.write(file.read())
            tmp_path = tmp.name
        text = docx2txt.process(tmp_path)
        os.unlink(tmp_path)
        return text.strip()
    except Exception:
        return ""

def extract_text_from_pptx(file):
    text = ""
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            tmp.write(file.read())
            tmp_path = tmp.name
        prs = pptx.Presentation(tmp_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + " "
        os.unlink(tmp_path)
    except Exception:
        return ""
    return text.strip()

def extract_text_from_txt(file):
    try:
        return file.read().decode("utf-8", errors="ignore").strip()
    except Exception:
        return ""

def extract_text(file, filetype):
    """
    filetype: extension with dot, e.g. '.pdf', '.docx', etc.
    Returns extracted text or empty string if failed.
    """
    if filetype == ".pdf":
        return extract_text_from_pdf(file)
    elif filetype == ".docx":
        return extract_text_from_docx(file)
    elif filetype == ".pptx":
        return extract_text_from_pptx(file)
    elif filetype == ".txt":
        return extract_text_from_txt(file)
    else:
        return "" 